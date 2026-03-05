"""
PPT 解析服务
============
使用 python-pptx 将 PPT 文件解析为纯文本
"""

from pptx import Presentation


def parse_ppt_to_text(file_path: str) -> str:
    """
    解析 PPT 文件，提取所有幻灯片中的文本内容

    Args:
        file_path: PPT 文件的绝对路径

    Returns:
        提取出的全部文本，按幻灯片分隔
    """
    prs = Presentation(file_path)
    all_text = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = []
        slide_texts.append(f"--- 第 {slide_num} 页 ---")

        for shape in slide.shapes:
            # 提取文本框中的文字
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_texts.append(text)

            # 提取表格中的文字
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text:
                            row_text.append(cell_text)
                    if row_text:
                        slide_texts.append(" | ".join(row_text))

        if len(slide_texts) > 1:  # 除了标题行外还有内容
            all_text.append("\n".join(slide_texts))

    return "\n\n".join(all_text)
